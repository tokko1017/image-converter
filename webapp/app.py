import streamlit as st
from PIL import Image
import pillow_heif
import io
import zipfile
import subprocess
import tempfile
import os
import csv
from pathlib import Path

pillow_heif.register_heif_opener()

st.set_page_config(
    page_title="写真・動画・文書 変換ツール",
    page_icon="🌸",
    layout="centered",
)

st.markdown("""
<style>
.stApp { background: linear-gradient(160deg, #fdf4ff 0%, #eff6ff 100%); }

/* ── タブ ── */
div[data-baseweb="tab-list"] {
    background: #f3e8ff;
    border-radius: 20px;
    padding: 6px;
    gap: 6px;
}
button[data-baseweb="tab"] {
    font-size: 1.05rem !important;
    font-weight: 700 !important;
    padding: 12px 28px !important;
    border-radius: 14px !important;
    color: #7c3aed !important;
    background: transparent !important;
}
button[data-baseweb="tab"]:hover {
    background: #ede9fe !important;
}
button[data-baseweb="tab"][aria-selected="true"] {
    background: linear-gradient(135deg, #e879f9, #818cf8) !important;
    color: white !important;
    box-shadow: 0 4px 15px rgba(168, 85, 247, 0.35) !important;
}
div[data-baseweb="tab-highlight"] { display: none !important; }
div[data-baseweb="tab-border"] { display: none !important; }

/* ── アップローダー ── */
div[data-testid="stFileUploader"] {
    border: 2.5px dashed #c084fc;
    border-radius: 16px;
    padding: 12px;
    background: #fdf4ff;
}

/* ── 変換ボタン ── */
div.stButton > button[kind="primary"] {
    background: linear-gradient(135deg, #e879f9, #818cf8);
    color: white; border: none;
    padding: 14px 0; font-size: 1.1rem;
    border-radius: 14px; width: 100%; font-weight: bold;
    box-shadow: 0 4px 15px rgba(168, 85, 247, 0.3);
    letter-spacing: 0.04em;
}
div.stButton > button[kind="primary"]:hover {
    background: linear-gradient(135deg, #f0abfc, #a5b4fc);
    box-shadow: 0 6px 20px rgba(168, 85, 247, 0.45);
}

/* ── ダウンロードボタン ── */
div[data-testid="stDownloadButton"] button {
    background: linear-gradient(135deg, #34d399, #059669);
    color: white; border: none;
    padding: 12px 0; font-size: 1rem;
    border-radius: 14px; width: 100%; font-weight: bold;
    box-shadow: 0 4px 12px rgba(52, 211, 153, 0.3);
}

/* ── 区切り線 ── */
hr {
    border: none; height: 2px;
    background: linear-gradient(90deg, #f9a8d4, #c084fc, #818cf8, #6ee7b7);
    border-radius: 2px; opacity: 0.55; margin: 16px 0;
}
</style>
""", unsafe_allow_html=True)

# ── ヘッダー ─────────────────────────────────────────────────────────────────
st.markdown("""
<div style="text-align:center; padding: 24px 0 16px;">
  <div style="font-size:2.8rem; margin-bottom:8px;">🌸</div>
  <div style="
    background: linear-gradient(135deg, #e879f9, #818cf8);
    -webkit-background-clip: text; -webkit-text-fill-color: transparent;
    background-clip: text;
    font-size: 2rem; font-weight: 900; margin: 0 0 6px; line-height: 1.3;">
    写真・動画・文書 変換ツール
  </div>
  <p style="color: #9333ea; font-size: 0.95rem; margin: 0;">
    ドラッグ＆ドロップで簡単変換 ✨
  </p>
</div>
""", unsafe_allow_html=True)

tab_img, tab_vid, tab_doc = st.tabs([
    "　🖼️  画像変換　",
    "　🎬  動画変換　",
    "　📄  文書変換　",
])

# ── ヘルパー ─────────────────────────────────────────────────────────────────
def step_guide(steps, formats_text):
    steps_html = "".join(f"""
      <div style="display:flex;align-items:center;gap:10px;padding:5px 0;">
        <span style="background:linear-gradient(135deg,#e879f9,#818cf8);color:white;
          border-radius:50%;width:26px;height:26px;display:flex;align-items:center;
          justify-content:center;font-size:0.8rem;font-weight:bold;flex-shrink:0;">{i+1}</span>
        <span style="color:#4c1d95;font-size:0.93rem;">{s}</span>
      </div>""" for i, s in enumerate(steps))
    st.markdown(f"""
    <div style="background:#fdf4ff;border-radius:16px;padding:20px;border:1.5px solid #e9d5ff;">
      <p style="color:#7c3aed;font-weight:700;font-size:0.95rem;margin:0 0 10px;">📋 使い方</p>
      {steps_html}
      <div style="border-top:1.5px solid #e9d5ff;margin-top:12px;padding-top:10px;">
        <p style="color:#7c3aed;font-size:0.85rem;margin:0;">
          <strong>対応形式：</strong>{formats_text}
        </p>
      </div>
    </div>""", unsafe_allow_html=True)

def format_badge(fmt, note):
    return (
        f"<div style='background:#f3e8ff;border-radius:10px;padding:8px 14px;"
        f"color:#6d28d9;font-size:0.9rem;margin-bottom:8px;'>"
        f"📌 <strong>{fmt}</strong>：{note}</div>"
    )

def file_badge(name, ext, size_label):
    return (
        f"<div style='padding:7px 12px;margin:4px 0;background:#fdf4ff;"
        f"border-radius:8px;border-left:3px solid #c084fc;font-size:0.9rem;'>"
        f"📄 <strong>{name}</strong> &nbsp;"
        f"<span style='background:#e9d5ff;color:#7c3aed;border-radius:4px;"
        f"padding:1px 7px;font-size:0.78rem;'>{ext}</span>"
        f" &nbsp; {size_label}</div>"
    )

def selected_badge(label):
    return (
        f"<div style='background:#f0fdf4;border-radius:12px;padding:12px 16px;"
        f"border:1.5px solid #6ee7b7;color:#065f46;font-weight:600;'>"
        f"✅ {label}</div>"
    )

def warning_badge(text):
    return (
        f"<div style='background:#fff7ed;border-radius:10px;padding:10px 14px;"
        f"border:1.5px solid #fed7aa;color:#92400e;font-size:0.85rem;margin-top:10px;'>"
        f"⚠️ {text}</div>"
    )

# ── 画像変換 ──────────────────────────────────────────────────────────────────
with tab_img:
    OUTPUT_FORMATS = {
        "JPEG": {"ext": ".jpg",  "pil": "JPEG", "alpha": False, "has_quality": True,  "note": "写真向け・高互換"},
        "PNG":  {"ext": ".png",  "pil": "PNG",  "alpha": True,  "has_quality": False, "note": "透過対応・可逆圧縮"},
        "WebP": {"ext": ".webp", "pil": "WEBP", "alpha": True,  "has_quality": True,  "note": "Web向け・高圧縮"},
        "AVIF": {"ext": ".avif", "pil": "AVIF", "alpha": True,  "has_quality": True,  "note": "次世代・高圧縮・透過対応"},
        "BMP":  {"ext": ".bmp",  "pil": "BMP",  "alpha": False, "has_quality": False, "note": "無圧縮・高互換性"},
        "TIFF": {"ext": ".tif",  "pil": "TIFF", "alpha": True,  "has_quality": False, "note": "印刷・業務向け"},
    }
    ACCEPT_TYPES = ["heic", "heif", "jpg", "jpeg", "png", "webp", "avif", "bmp", "tiff", "tif"]

    st.markdown("<div style='height:10px;'></div>", unsafe_allow_html=True)

    col1, col2 = st.columns([3, 2])
    with col1:
        fmt = st.selectbox("✨ 変換後の形式", list(OUTPUT_FORMATS.keys()),
                           help="アップロードした画像をこの形式に変換します")
    info = OUTPUT_FORMATS[fmt]
    st.markdown(format_badge(fmt, info["note"]), unsafe_allow_html=True)

    with col2:
        if info["has_quality"]:
            quality = st.slider("🎚️ 画質", min_value=40, max_value=100, value=85,
                                help="数値が高いほど高画質・ファイルサイズ大")
        else:
            quality = 85
            st.info("可逆圧縮のため画質設定なし", icon="ℹ️")

    st.divider()

    uploaded_files = st.file_uploader(
        "📂 画像をアップロード（複数まとめてOK）",
        type=ACCEPT_TYPES,
        accept_multiple_files=True,
    )

    if not uploaded_files:
        step_guide(
            ["画像ファイルを上の枠にドラッグ＆ドロップ（または「Upload」をクリック）",
             "変換後の形式・画質を選ぶ",
             "「変換開始」ボタンを押す",
             "ダウンロードボタンから保存する"],
            "HEIC / HEIF / JPEG / PNG / WebP / AVIF / BMP / TIFF"
        )
    else:
        st.markdown(selected_badge(f"{len(uploaded_files)} ファイルを選択中"), unsafe_allow_html=True)
        st.markdown("<div style='height:8px;'></div>", unsafe_allow_html=True)
        for f in uploaded_files:
            size_kb = round(f.size / 1024, 1)
            ext = Path(f.name).suffix.upper().lstrip(".")
            st.markdown(file_badge(f.name, ext, f"{size_kb} KB"), unsafe_allow_html=True)
        st.divider()

        if st.button("✨ 変換開始", type="primary", key="img_convert"):
            results = []
            errors = []
            progress_bar = st.progress(0, text="変換中...")
            total = len(uploaded_files)

            for i, file in enumerate(uploaded_files):
                progress_bar.progress((i + 1) / total, text=f"変換中... {file.name} ({i+1}/{total})")
                try:
                    img = Image.open(file)
                    if not info["alpha"] and img.mode in ("RGBA", "LA", "P"):
                        bg = Image.new("RGB", img.size, (255, 255, 255))
                        src = img.convert("RGBA") if img.mode != "RGBA" else img
                        bg.paste(src, mask=src.split()[3])
                        img_out = bg
                    elif info["alpha"] and img.mode == "P":
                        img_out = img.convert("RGBA")
                    else:
                        img_out = img

                    if not info["alpha"] and img_out.mode not in ("RGB", "L"):
                        img_out = img_out.convert("RGB")
                    elif info["alpha"] and img_out.mode not in ("RGB", "RGBA", "L", "LA"):
                        img_out = img_out.convert("RGBA")

                    buf = io.BytesIO()
                    save_kwargs = {}
                    if info["has_quality"]:
                        save_kwargs["quality"] = quality
                    if fmt == "JPEG":
                        save_kwargs["optimize"] = True
                    img_out.save(buf, info["pil"], **save_kwargs)
                    buf.seek(0)
                    results.append((Path(file.name).stem + info["ext"], buf.getvalue()))
                except Exception as e:
                    errors.append(f"{file.name}：{e}")

            progress_bar.empty()
            for msg in errors:
                st.error(f"失敗：{msg}")

            if results:
                st.success(f"🎉 {len(results)} ファイルの変換が完了しました！")
                if len(results) == 1:
                    name, data = results[0]
                    st.download_button(f"⬇️ ダウンロード：{name}", data=data, file_name=name)
                else:
                    zip_buf = io.BytesIO()
                    with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:
                        for name, data in results:
                            zf.writestr(name, data)
                    zip_buf.seek(0)
                    st.download_button(
                        f"⬇️ すべてダウンロード（{len(results)} ファイル） .zip",
                        data=zip_buf.getvalue(),
                        file_name="converted_images.zip",
                        mime="application/zip",
                    )

# ── 動画変換 ──────────────────────────────────────────────────────────────────
with tab_vid:
    VIDEO_INPUT_TYPES = ["mp4", "mov", "avi", "mkv", "webm", "m4v"]
    VIDEO_OUTPUT_FORMATS = {
        "MP4":  {"ext": ".mp4",  "note": "最も互換性が高い・スマホ・PC問わず再生可能"},
        "WebM": {"ext": ".webm", "note": "Web向け・軽量・ブラウザ再生に最適"},
        "GIF":  {"ext": ".gif",  "note": "アニメーションGIF・SNS向け（音声なし）"},
        "MP3":  {"ext": ".mp3",  "note": "音声のみ抽出して保存"},
    }

    def ffmpeg_available():
        try:
            return subprocess.run(["ffmpeg", "-version"], capture_output=True, timeout=5).returncode == 0
        except (FileNotFoundError, subprocess.TimeoutExpired):
            return False

    st.markdown("<div style='height:10px;'></div>", unsafe_allow_html=True)

    if not ffmpeg_available():
        st.error("FFmpegがインストールされていないため、動画変換は使用できません。")
    else:
        col1v, col2v = st.columns([3, 2])
        with col1v:
            vfmt = st.selectbox("✨ 変換後の形式", list(VIDEO_OUTPUT_FORMATS.keys()),
                                help="動画をこの形式に変換します", key="vfmt")
        vinfo = VIDEO_OUTPUT_FORMATS[vfmt]
        st.markdown(format_badge(vfmt, vinfo["note"]), unsafe_allow_html=True)

        gif_fps = 10
        gif_width = 480
        if vfmt == "GIF":
            with col2v:
                gif_fps = st.select_slider("🎞️ フレームレート (fps)", options=[5, 10, 15, 24], value=10)
                gif_width = st.select_slider("📐 横幅 (px)", options=[240, 360, 480, 640], value=480)

        st.divider()

        uploaded_video = st.file_uploader(
            "📂 動画をアップロード（1ファイルずつ）",
            type=VIDEO_INPUT_TYPES,
            accept_multiple_files=False,
            key="video_uploader",
        )

        if not uploaded_video:
            step_guide(
                ["動画ファイルを上の枠にドラッグ＆ドロップ（または「Upload」をクリック）",
                 "変換後の形式を選ぶ（GIFはフレームレート・横幅も設定可）",
                 "「変換開始」ボタンを押す",
                 "ダウンロードボタンから保存する"],
                "MP4 / MOV / AVI / MKV / WebM / M4V"
            )
            st.markdown("""
            <div style="background:#f3e8ff;border-radius:14px;padding:16px;margin-top:12px;border:1.5px solid #ddd6fe;">
              <p style="color:#7c3aed;font-weight:700;margin:0 0 10px;">📊 変換できる組み合わせ</p>
              <p style="color:#6d28d9;font-size:0.85rem;margin:0 0 8px;">すべての入力形式から、すべての出力形式へ変換できます。</p>
              <table style="width:100%;border-collapse:collapse;font-size:0.88rem;color:#4c1d95;">
                <tr style="background:#ede9fe;">
                  <th style="padding:6px 10px;text-align:left;border-radius:6px 0 0 6px;">出力形式</th>
                  <th style="padding:6px 10px;text-align:left;border-radius:0 6px 6px 0;">説明</th>
                </tr>
                <tr><td style="padding:6px 10px;">MP4</td><td style="padding:6px 10px;">最も互換性が高い・スマホ・PC問わず再生可能</td></tr>
                <tr style="background:#faf5ff;"><td style="padding:6px 10px;">WebM</td><td style="padding:6px 10px;">Web向け・軽量・ブラウザ再生に最適</td></tr>
                <tr><td style="padding:6px 10px;">GIF</td><td style="padding:6px 10px;">アニメーションGIF・SNS向け（音声なし）</td></tr>
                <tr style="background:#faf5ff;"><td style="padding:6px 10px;">MP3</td><td style="padding:6px 10px;">音声のみ抽出して保存</td></tr>
              </table>
            </div>
            """, unsafe_allow_html=True)
            st.markdown(
                warning_badge("ファイルサイズは <strong>200MB以下</strong> を推奨します"),
                unsafe_allow_html=True
            )
        else:
            size_mb = round(uploaded_video.size / 1024 / 1024, 1)
            st.markdown(
                selected_badge(f"{uploaded_video.name} &nbsp; <span style='font-weight:400;'>{size_mb} MB</span>"),
                unsafe_allow_html=True
            )
            st.divider()

            if st.button("✨ 変換開始", type="primary", key="vid_convert"):
                with st.spinner("変換中... しばらくお待ちください 🎬"):
                    try:
                        with tempfile.TemporaryDirectory() as tmpdir:
                            input_path = os.path.join(tmpdir, "input" + Path(uploaded_video.name).suffix)
                            with open(input_path, "wb") as f:
                                f.write(uploaded_video.getbuffer())

                            out_name = Path(uploaded_video.name).stem + vinfo["ext"]
                            output_path = os.path.join(tmpdir, out_name)

                            if vfmt == "MP4":
                                cmd = ["ffmpeg", "-i", input_path,
                                       "-c:v", "libx264", "-c:a", "aac",
                                       "-movflags", "+faststart", "-y", output_path]
                            elif vfmt == "WebM":
                                cmd = ["ffmpeg", "-i", input_path,
                                       "-c:v", "libvpx-vp9", "-c:a", "libopus",
                                       "-y", output_path]
                            elif vfmt == "GIF":
                                palette_path = os.path.join(tmpdir, "palette.png")
                                subprocess.run(
                                    ["ffmpeg", "-i", input_path,
                                     "-vf", f"fps={gif_fps},scale={gif_width}:-1:flags=lanczos,palettegen",
                                     "-y", palette_path],
                                    capture_output=True, timeout=120,
                                )
                                cmd = ["ffmpeg", "-i", input_path, "-i", palette_path,
                                       "-filter_complex",
                                       f"fps={gif_fps},scale={gif_width}:-1:flags=lanczos[x];[x][1:v]paletteuse",
                                       "-y", output_path]
                            elif vfmt == "MP3":
                                cmd = ["ffmpeg", "-i", input_path,
                                       "-q:a", "2", "-map", "a", "-y", output_path]

                            result = subprocess.run(cmd, capture_output=True, timeout=300)

                            if result.returncode != 0:
                                err_msg = result.stderr.decode("utf-8", errors="replace")
                                st.error(f"変換に失敗しました。\n```\n{err_msg[-500:]}\n```")
                            else:
                                with open(output_path, "rb") as f:
                                    output_data = f.read()
                                out_size_mb = round(len(output_data) / 1024 / 1024, 1)
                                st.success("🎉 変換が完了しました！")
                                st.caption(f"出力サイズ：{out_size_mb} MB")
                                st.download_button(
                                    label=f"⬇️ ダウンロード：{out_name}",
                                    data=output_data,
                                    file_name=out_name,
                                )

                    except subprocess.TimeoutExpired:
                        st.error("⏱️ 変換がタイムアウトしました。ファイルサイズを小さくしてお試しください。")
                    except Exception as e:
                        st.error(f"エラーが発生しました：{e}")

# ── 文書変換 ──────────────────────────────────────────────────────────────────
with tab_doc:
    DOC_INPUT_TYPES = ["txt", "docx", "pdf", "xlsx", "pptx"]

    OUTPUT_BY_EXT = {
        "txt":  ["PDF", "JPEG", "PNG", "WebP"],
        "docx": ["PDF", "TXT", "JPEG", "PNG", "WebP"],
        "pdf":  ["TXT", "JPEG", "PNG", "WebP"],
        "xlsx": ["PDF", "CSV", "JPEG", "PNG", "WebP"],
        "pptx": ["PDF", "TXT", "JPEG", "PNG", "WebP"],
    }

    FORMAT_NOTES = {
        "PDF":  "どの端末でも同じレイアウトで開ける",
        "TXT":  "テキストを抽出してプレーンテキストで保存",
        "CSV":  "Excel・スプレッドシートで開けるデータ形式",
        "JPEG": "ページ・スライドごとに画像化（写真向け・高互換）",
        "PNG":  "ページ・スライドごとに画像化（高品質・透過対応）",
        "WebP": "ページ・スライドごとに画像化（Web向け・軽量）",
    }

    def libreoffice_available():
        try:
            return subprocess.run(["libreoffice", "--version"],
                                  capture_output=True, timeout=5).returncode == 0
        except (FileNotFoundError, subprocess.TimeoutExpired):
            return False

    def to_pdf(input_path, tmpdir):
        env = {**os.environ, "HOME": tmpdir}
        cmd = ["libreoffice", "--headless", "--norestore",
               "--convert-to", "pdf", "--outdir", tmpdir, input_path]
        result = subprocess.run(cmd, capture_output=True, timeout=120, env=env)
        out = Path(tmpdir) / (Path(input_path).stem + ".pdf")
        if result.returncode == 0 and out.exists():
            return out, None
        return None, result.stderr.decode("utf-8", errors="replace")

    def extract_text(input_path):
        ext = Path(input_path).suffix.lower()
        if ext == ".pdf":
            import pdfplumber
            parts = []
            with pdfplumber.open(input_path) as pdf:
                for page in pdf.pages:
                    t = page.extract_text()
                    if t:
                        parts.append(t)
            return "\n\n".join(parts)
        elif ext == ".docx":
            from docx import Document
            doc = Document(input_path)
            return "\n".join(p.text for p in doc.paragraphs)
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(input_path)
            lines = []
            for i, slide in enumerate(prs.slides, 1):
                lines.append(f"=== スライド {i} ===")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        lines.append(shape.text_frame.text)
                lines.append("")
            return "\n".join(lines)
        return ""

    def to_images(input_path, img_format, dpi, tmpdir):
        import fitz
        ext = Path(input_path).suffix.lower()
        if ext != ".pdf":
            pdf_path, err = to_pdf(input_path, tmpdir)
            if not pdf_path:
                return None, err
            input_path = str(pdf_path)
        doc = fitz.open(input_path)
        matrix = fitz.Matrix(dpi / 72, dpi / 72)
        suffix = {"JPEG": ".jpg", "PNG": ".png", "WebP": ".webp"}[img_format]
        pil_fmt = {"JPEG": "JPEG", "PNG": "PNG", "WebP": "WEBP"}[img_format]
        images = []
        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=matrix, alpha=False)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            buf = io.BytesIO()
            if img_format == "JPEG":
                img.save(buf, pil_fmt, quality=90, optimize=True)
            elif img_format == "PNG":
                img.save(buf, pil_fmt)
            else:
                img.save(buf, pil_fmt, quality=90)
            buf.seek(0)
            images.append((f"page_{i+1:03d}{suffix}", buf.getvalue()))
        doc.close()
        return images, None

    def to_csv_bytes(input_path):
        import openpyxl
        wb = openpyxl.load_workbook(input_path, data_only=True)
        ws = wb.active
        buf = io.StringIO()
        writer = csv.writer(buf)
        for row in ws.iter_rows(values_only=True):
            writer.writerow(['' if v is None else v for v in row])
        return buf.getvalue().encode("utf-8-sig")

    st.markdown("<div style='height:10px;'></div>", unsafe_allow_html=True)

    uploaded_doc = st.file_uploader(
        "📂 文書をアップロード",
        type=DOC_INPUT_TYPES,
        accept_multiple_files=False,
        key="doc_uploader",
    )

    if not uploaded_doc:
        step_guide(
            ["文書ファイルを上の枠にドラッグ＆ドロップ（または「Upload」をクリック）",
             "変換後の形式を選ぶ（ファイルの種類によって選択肢が変わります）",
             "「変換開始」ボタンを押す",
             "ダウンロードボタンから保存する"],
            "TXT / DOCX / PDF / XLSX / PPTX"
        )
        st.markdown("""
        <div style="background:#f3e8ff;border-radius:14px;padding:16px;margin-top:12px;border:1.5px solid #ddd6fe;">
          <p style="color:#7c3aed;font-weight:700;margin:0 0 10px;">📊 変換できる組み合わせ</p>
          <table style="width:100%;border-collapse:collapse;font-size:0.88rem;color:#4c1d95;">
            <tr style="background:#ede9fe;">
              <th style="padding:6px 10px;text-align:left;border-radius:6px 0 0 6px;">入力</th>
              <th style="padding:6px 10px;text-align:left;border-radius:0 6px 6px 0;">変換できる形式</th>
            </tr>
            <tr><td style="padding:6px 10px;">TXT</td><td style="padding:6px 10px;">→ PDF・JPEG・PNG・WebP</td></tr>
            <tr style="background:#faf5ff;"><td style="padding:6px 10px;">DOCX（Word）</td><td style="padding:6px 10px;">→ PDF・TXT・JPEG・PNG・WebP</td></tr>
            <tr><td style="padding:6px 10px;">PDF</td><td style="padding:6px 10px;">→ TXT・JPEG・PNG・WebP</td></tr>
            <tr style="background:#faf5ff;"><td style="padding:6px 10px;">XLSX（Excel）</td><td style="padding:6px 10px;">→ PDF・CSV・JPEG・PNG・WebP</td></tr>
            <tr><td style="padding:6px 10px;">PPTX（PowerPoint）</td><td style="padding:6px 10px;">→ PDF・TXT・JPEG・PNG・WebP</td></tr>
          </table>
        </div>
        """, unsafe_allow_html=True)
    else:
        ext = Path(uploaded_doc.name).suffix.lower().lstrip(".")
        size_kb = round(uploaded_doc.size / 1024, 1)
        st.markdown(
            selected_badge(f"{uploaded_doc.name} &nbsp; <span style='font-weight:400;'>{size_kb} KB</span>"),
            unsafe_allow_html=True
        )

        available_outputs = OUTPUT_BY_EXT.get(ext, [])
        if not available_outputs:
            st.error("対応していない形式です。")
        else:
            st.markdown("<div style='height:8px;'></div>", unsafe_allow_html=True)
            out_fmt = st.selectbox("✨ 変換後の形式", available_outputs, key="doc_out_fmt")
            st.markdown(format_badge(out_fmt, FORMAT_NOTES[out_fmt]), unsafe_allow_html=True)

            dpi = 150
            if out_fmt in ("JPEG", "PNG", "WebP"):
                dpi = st.select_slider(
                    "🖨️ 解像度 (DPI)",
                    options=[96, 150, 300],
                    value=150,
                    help="96=画面表示向け・150=標準・300=印刷品質（ファイルサイズ大）",
                    key="doc_dpi",
                )
            st.divider()

            if st.button("✨ 変換開始", type="primary", key="doc_convert"):
                needs_lo = (out_fmt in ("PDF", "JPEG", "PNG", "WebP") and ext != "pdf")
                if needs_lo and not libreoffice_available():
                    st.error("LibreOfficeがインストールされていないため、PDF変換は使用できません。")
                else:
                    with st.spinner("変換中... しばらくお待ちください 📄"):
                        try:
                            with tempfile.TemporaryDirectory() as tmpdir:
                                input_path = os.path.join(tmpdir, "input." + ext)
                                with open(input_path, "wb") as f:
                                    f.write(uploaded_doc.getbuffer())

                                out_stem = Path(uploaded_doc.name).stem

                                if out_fmt == "PDF":
                                    out_path, err = to_pdf(input_path, tmpdir)
                                    if out_path:
                                        data = out_path.read_bytes()
                                        st.success("🎉 変換が完了しました！")
                                        st.download_button(
                                            f"⬇️ ダウンロード：{out_stem}.pdf",
                                            data=data,
                                            file_name=out_stem + ".pdf",
                                            mime="application/pdf",
                                        )
                                    else:
                                        st.error(f"変換に失敗しました。\n```\n{err[-500:]}\n```")

                                elif out_fmt == "TXT":
                                    text = extract_text(input_path)
                                    if text.strip():
                                        st.success("🎉 テキストの抽出が完了しました！")
                                        st.download_button(
                                            f"⬇️ ダウンロード：{out_stem}.txt",
                                            data=text.encode("utf-8"),
                                            file_name=out_stem + ".txt",
                                            mime="text/plain",
                                        )
                                    else:
                                        st.warning("テキストが抽出できませんでした。スキャンされたPDFなどは対応できません。")

                                elif out_fmt == "CSV":
                                    data = to_csv_bytes(input_path)
                                    st.success("🎉 変換が完了しました！")
                                    st.download_button(
                                        f"⬇️ ダウンロード：{out_stem}.csv",
                                        data=data,
                                        file_name=out_stem + ".csv",
                                        mime="text/csv",
                                    )

                                elif out_fmt in ("JPEG", "PNG", "WebP"):
                                    images, err = to_images(input_path, out_fmt, dpi, tmpdir)
                                    if images:
                                        st.success(f"🎉 {len(images)} ページの画像化が完了しました！")
                                        if len(images) == 1:
                                            name, data = images[0]
                                            st.download_button(
                                                f"⬇️ ダウンロード：{name}",
                                                data=data,
                                                file_name=name,
                                            )
                                        else:
                                            zip_buf = io.BytesIO()
                                            with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:
                                                for name, data in images:
                                                    zf.writestr(f"{out_stem}/{name}", data)
                                            zip_buf.seek(0)
                                            st.download_button(
                                                f"⬇️ すべてダウンロード（{len(images)} ページ） .zip",
                                                data=zip_buf.getvalue(),
                                                file_name=f"{out_stem}_images.zip",
                                                mime="application/zip",
                                            )
                                    else:
                                        st.error(f"画像変換に失敗しました。\n```\n{err[-500:]}\n```")

                        except subprocess.TimeoutExpired:
                            st.error("⏱️ 変換がタイムアウトしました。")
                        except Exception as e:
                            st.error(f"エラーが発生しました：{e}")
