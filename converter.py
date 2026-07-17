import tkinter as tk
from tkinter import ttk, filedialog, messagebox
import threading
import os
import io
import csv
import zipfile
import subprocess
import tempfile
import sys
from pathlib import Path
from PIL import Image
import pillow_heif

pillow_heif.register_heif_opener()

_NO_WIN = subprocess.CREATE_NO_WINDOW if sys.platform == "win32" else 0

# ── カラーパレット（薄水色ベース・紺青統一） ─────────────────────────────────
BG      = "#f0f7ff"   # 薄い水色背景
SURFACE = "#dbeafe"   # カード背景（少し濃い水色）
CARD    = "#f0f7ff"
BORDER  = "#93c5fd"   # ボーダー
BORDER2 = "#bfdbfe"
TAB_BAR = "#dbeafe"
TAB_ACT = "#1e3a8a"   # アクティブタブ（紺）
TAB_HOV = "#bfdbfe"
BTN     = "#1e3a8a"   # プライマリボタン（紺）
BTN_HOV = "#1e40af"
TEXT    = "#1e3a8a"   # 濃いテキスト（紺）
MID     = "#1d4ed8"   # 中間テキスト（青）
DIM     = "#60a5fa"   # 薄テキスト（明るい青）
ROW_ALT = "#eff6ff"   # リスト交互行
SUCCESS = "#059669"
ERROR_C = "#dc2626"
WARN_C  = "#d97706"

# ── ユーティリティ ────────────────────────────────────────────────────────────
def ffmpeg_ok():
    try:
        return subprocess.run(["ffmpeg", "-version"],
                              capture_output=True, timeout=5,
                              creationflags=_NO_WIN).returncode == 0
    except Exception:
        return False

def libreoffice_ok():
    for cmd in ["libreoffice", "soffice"]:
        try:
            r = subprocess.run([cmd, "--version"], capture_output=True,
                               timeout=5, creationflags=_NO_WIN)
            if r.returncode == 0:
                return cmd
        except Exception:
            pass
    return None

def tesseract_path():
    import shutil as _shutil
    exe = _shutil.which("tesseract")
    if exe:
        return exe
    default = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
    if os.path.isfile(default):
        return default
    return None

def tesseract_ok():
    return tesseract_path() is not None

def bundled_tessdata_dir():
    base = getattr(sys, "_MEIPASS", None) or os.path.dirname(os.path.abspath(__file__))
    d = os.path.join(base, "tessdata")
    return d if os.path.isdir(d) else None

def tesseract_langs():
    langs = set()
    exe = tesseract_path()
    if exe:
        try:
            r = subprocess.run([exe, "--list-langs"], capture_output=True,
                               timeout=5, creationflags=_NO_WIN)
            out = (r.stdout + r.stderr).decode("utf-8", errors="ignore")
            langs |= {ln.strip() for ln in out.splitlines()[1:] if ln.strip()}
        except Exception:
            pass
    bundled = bundled_tessdata_dir()
    if bundled:
        for fn in os.listdir(bundled):
            if fn.endswith(".traineddata"):
                langs.add(fn[:-len(".traineddata")])
    return langs

def _run_tesseract(exe, img, lang, tessdata_dir=None, psm=None, want_conf=False, timeout=60):
    """tesseract.exeを直接subprocessで呼び出す（パスにスペースや日本語を含んでいても安全）。"""
    with tempfile.TemporaryDirectory() as td:
        img_path = os.path.join(td, "page.png")
        img.save(img_path, "PNG")
        out_base = os.path.join(td, "out")
        cmd = [exe, img_path, out_base, "-l", lang]
        if psm is not None:
            cmd += ["--psm", str(psm)]
        if tessdata_dir:
            cmd += ["--tessdata-dir", tessdata_dir]
        if want_conf:
            cmd += ["-c", "tessedit_create_tsv=1"]
        try:
            subprocess.run(cmd, capture_output=True, timeout=timeout, creationflags=_NO_WIN)
        except Exception:
            return (-1.0 if want_conf else "")
        if want_conf:
            tsv_path = out_base + ".tsv"
            if not os.path.isfile(tsv_path):
                return -1.0
            vals = []
            with open(tsv_path, encoding="utf-8", errors="replace") as f:
                reader = csv.DictReader(f, delimiter="\t")
                for row in reader:
                    try:
                        c = float(row.get("conf", -1))
                        if c >= 0:
                            vals.append(c)
                    except (TypeError, ValueError):
                        pass
            return sum(vals) / len(vals) if vals else -1.0
        else:
            txt_path = out_base + ".txt"
            if os.path.isfile(txt_path):
                return Path(txt_path).read_text(encoding="utf-8", errors="replace")
            return ""

import re
_CJK_GAP_RE = re.compile(r'(?<=[ぁ-んァ-ヶ一-龥々〆〤、。「」！？])\s+(?=[ぁ-んァ-ヶ一-龥々〆〤、。「」！？])')

def clean_ocr_ja(text):
    """Tesseractが日本語OCR時に文字間へ挿入する余分な空白を除去する。"""
    return _CJK_GAP_RE.sub("", text)

# ── ピル型ボタン（Canvas製） ──────────────────────────────────────────────────
class PillBtn(tk.Canvas):
    def __init__(self, parent, text, cmd, height=46):
        super().__init__(parent, height=height, bg=parent.cget("bg"),
                         highlightthickness=0, cursor="hand2")
        self._c  = BTN
        self._h  = BTN_HOV
        self._tx = text
        self._cm = cmd
        self.bind("<Configure>", self._draw)
        self.bind("<Button-1>",  lambda e: cmd())
        self.bind("<Enter>",     lambda e: self._fill(self._h))
        self.bind("<Leave>",     lambda e: self._fill(BTN))

    def _draw(self, e):
        self.delete("all")
        w, h = e.width, e.height
        r = h // 2
        self.create_oval(0, 0, h, h, fill=BTN, outline="")
        self.create_oval(w - h, 0, w, h, fill=BTN, outline="")
        self.create_rectangle(r, 0, w - r, h, fill=BTN, outline="")
        self.create_text(w // 2, h // 2, text=self._tx,
                         fill="#ffffff", font=("Yu Gothic UI", 12, "bold"))

    def _fill(self, color):
        for item in self.find_all():
            t = self.type(item)
            if t in ("oval", "rectangle"):
                self.itemconfig(item, fill=color)

# ── セカンダリボタン ──────────────────────────────────────────────────────────
class SecBtn(tk.Button):
    def __init__(self, parent, text, cmd, **kw):
        super().__init__(parent, text=text, command=cmd,
                         bg=SURFACE, fg=MID,
                         activebackground=TAB_HOV, activeforeground=MID,
                         font=("Yu Gothic UI", 9),
                         relief="flat", bd=0, cursor="hand2",
                         padx=12, pady=5,
                         highlightbackground=BORDER, highlightthickness=1,
                         **kw)
        self.bind("<Enter>", lambda _: self.config(bg=TAB_HOV))
        self.bind("<Leave>", lambda _: self.config(bg=SURFACE))

# ── 縁取りカード ──────────────────────────────────────────────────────────────
def make_card(parent):
    outer = tk.Frame(parent, bg=BORDER2,
                     highlightbackground=BORDER2, highlightthickness=0)
    inner = tk.Frame(outer, bg=SURFACE)
    inner.pack(fill="both", expand=True, padx=2, pady=2)
    return outer, inner

# ── カラープログレスバー ──────────────────────────────────────────────────────
class ColorBar(tk.Frame):
    def __init__(self, parent, var, height=10):
        super().__init__(parent, bg=BORDER2, height=height,
                         highlightbackground=BORDER2, highlightthickness=1)
        self._bar = tk.Frame(self, bg=BTN, height=height)
        self._bar.place(x=0, y=0, relheight=1, relwidth=0)
        var.trace_add("write", lambda *_: self._bar.place(
            relwidth=min(var.get() / 100.0, 1.0)))

# ── ガイドカード ──────────────────────────────────────────────────────────────
def build_guide(parent, steps, formats):
    outer, inner = make_card(parent)
    outer.pack(fill="x", padx=16, pady=(14, 0))
    tk.Label(inner, text="📋 使い方",
             font=("Yu Gothic UI", 10, "bold"),
             bg=SURFACE, fg=MID,
             anchor="w", padx=12, pady=6).pack(fill="x")
    for s in steps:
        tk.Label(inner, text=s, font=("Yu Gothic UI", 9),
                 bg=SURFACE, fg=TEXT,
                 anchor="w", padx=20, pady=2).pack(fill="x")
    tk.Frame(inner, bg=BORDER, height=1).pack(fill="x", padx=12, pady=(6, 0))
    tk.Label(inner, text=f"対応形式：{formats}",
             font=("Yu Gothic UI", 8),
             bg=SURFACE, fg=MID,
             anchor="w", padx=12, pady=6).pack(fill="x")

# ── コンボボックス ────────────────────────────────────────────────────────────
def build_cb(parent, var, values):
    style = ttk.Style()
    style.configure("P.TCombobox",
                    fieldbackground=CARD, background=CARD,
                    foreground=TEXT, selectbackground=SURFACE,
                    arrowcolor=MID)
    return ttk.Combobox(parent, textvariable=var, values=values,
                        state="readonly", width=9,
                        style="P.TCombobox",
                        font=("Yu Gothic UI", 10))


# ════════════════════════════════════════════════════════════════════════════
# 画像変換タブ
# ════════════════════════════════════════════════════════════════════════════
class ImageTab(tk.Frame):
    IMG_EXTS = {".heic", ".heif", ".jpg", ".jpeg", ".png",
                ".webp", ".avif", ".bmp", ".tiff", ".tif"}
    FORMATS = {
        "JPEG": {"ext": ".jpg",  "pil": "JPEG", "alpha": False, "quality": True,  "note": "写真向け・高互換"},
        "PNG":  {"ext": ".png",  "pil": "PNG",  "alpha": True,  "quality": False, "note": "透過対応・可逆圧縮"},
        "WebP": {"ext": ".webp", "pil": "WEBP", "alpha": True,  "quality": True,  "note": "Web向け・高圧縮"},
        "AVIF": {"ext": ".avif", "pil": "AVIF", "alpha": True,  "quality": True,  "note": "次世代・超高圧縮"},
        "BMP":  {"ext": ".bmp",  "pil": "BMP",  "alpha": False, "quality": False, "note": "無圧縮・高互換"},
        "TIFF": {"ext": ".tif",  "pil": "TIFF", "alpha": True,  "quality": False, "note": "印刷・業務向け"},
    }
    FILETYPES = [("対応画像", "*.heic *.heif *.jpg *.jpeg *.png *.webp *.avif *.bmp *.tiff *.tif"), ("すべて", "*.*")]

    def __init__(self, parent):
        super().__init__(parent, bg=BG)
        self.files = []
        self.converting = False
        self.out_format = tk.StringVar(value="JPEG")
        self.quality    = tk.IntVar(value=85)
        self.output_dir = tk.StringVar(value="")
        self._build()

    def _build(self):
        build_guide(self,
            ["① ファイル追加 または フォルダ追加 でファイルを選ぶ",
             "② 変換後の形式・画質を設定する",
             "③「✨ 変換開始」ボタンを押す"],
            "HEIC / HEIF / JPEG / PNG / WebP / AVIF / BMP / TIFF")

        btn_row = tk.Frame(self, bg=BG)
        btn_row.pack(fill="x", padx=16, pady=(10, 0))
        SecBtn(btn_row, "📂 ファイル追加",   self._add_files).pack(side="left", padx=(0, 6))
        SecBtn(btn_row, "📁 フォルダ追加",   self._add_folder).pack(side="left")
        SecBtn(btn_row, "🗑 クリア", self._clear).pack(side="right")

        lo, li = make_card(self)
        lo.pack(fill="both", expand=True, padx=16, pady=8)
        self.canvas = tk.Canvas(li, bg=SURFACE, highlightthickness=0)
        sb = ttk.Scrollbar(li, orient="vertical", command=self.canvas.yview)
        self.canvas.configure(yscrollcommand=sb.set)
        sb.pack(side="right", fill="y")
        self.canvas.pack(side="left", fill="both", expand=True)
        self.inner = tk.Frame(self.canvas, bg=SURFACE)
        self._cwin = self.canvas.create_window((0, 0), window=self.inner, anchor="nw")
        self.inner.bind("<Configure>",
                        lambda e: self.canvas.configure(scrollregion=self.canvas.bbox("all")))
        self.canvas.bind("<Configure>",
                         lambda e: self.canvas.itemconfig(self._cwin, width=e.width))
        self._show_empty()

        co, ci = make_card(self)
        co.pack(fill="x", padx=16, pady=(0, 6))

        r1 = tk.Frame(ci, bg=SURFACE)
        r1.pack(fill="x", padx=12, pady=(10, 4))
        tk.Label(r1, text="✨ 変換後の形式", font=("Yu Gothic UI", 10, "bold"),
                 bg=SURFACE, fg=MID).pack(side="left")
        cb = build_cb(r1, self.out_format, list(self.FORMATS.keys()))
        cb.pack(side="left", padx=(8, 0))
        cb.bind("<<ComboboxSelected>>", self._on_fmt)
        self.fmt_note = tk.Label(r1, text="", font=("Yu Gothic UI", 9),
                                 bg=SURFACE, fg=MID)
        self.fmt_note.pack(side="left", padx=(10, 0))

        self.q_frame = tk.Frame(ci, bg=SURFACE)
        self.q_frame.pack(fill="x", padx=12, pady=(0, 4))
        tk.Label(self.q_frame, text="🎚️ 画質",
                 font=("Yu Gothic UI", 9), bg=SURFACE, fg=DIM,
                 width=7, anchor="w").pack(side="left")
        style = ttk.Style()
        style.configure("P.Horizontal.TScale", background=SURFACE)
        ttk.Scale(self.q_frame, from_=40, to=100, variable=self.quality,
                  orient="horizontal", style="P.Horizontal.TScale",
                  command=lambda _: self.q_val.config(text=f"{self.quality.get()}%")).pack(
            side="left", fill="x", expand=True, padx=(4, 8))
        self.q_val = tk.Label(self.q_frame, text="85%",
                              font=("Yu Gothic UI", 10, "bold"),
                              bg=SURFACE, fg=BTN, width=5)
        self.q_val.pack(side="right")

        r3 = tk.Frame(ci, bg=SURFACE)
        r3.pack(fill="x", padx=12, pady=(0, 10))
        tk.Label(r3, text="出力先", font=("Yu Gothic UI", 9),
                 bg=SURFACE, fg=DIM, width=7, anchor="w").pack(side="left")
        tk.Entry(r3, textvariable=self.output_dir,
                 bg=CARD, fg=TEXT, insertbackground=MID,
                 relief="flat", font=("Yu Gothic UI", 9), bd=3,
                 highlightbackground=BORDER, highlightthickness=1).pack(
            side="left", fill="x", expand=True, padx=(4, 8))
        SecBtn(r3, "選択", self._choose_out).pack(side="left")

        self.prog_label = tk.Label(self, text="", font=("Yu Gothic UI", 9),
                                   bg=BG, fg=MID)
        self.prog_label.pack(padx=16, anchor="w")
        self.prog_var = tk.DoubleVar(value=0)
        ColorBar(self, self.prog_var).pack(fill="x", padx=16, pady=(2, 8))

        PillBtn(self, "✨ 変換開始", self._start).pack(
            fill="x", padx=16, pady=(0, 16))
        self._on_fmt()

    def _on_fmt(self, _=None):
        fmt = self.out_format.get()
        has_q = self.FORMATS[fmt]["quality"]
        self.fmt_note.config(text=f"📌 {self.FORMATS[fmt]['note']}")
        state = "normal" if has_q else "disabled"
        for w in self.q_frame.winfo_children():
            try:
                w.configure(state=state)
            except Exception:
                pass

    def _add_files(self):
        paths = filedialog.askopenfilenames(title="画像を選択", filetypes=self.FILETYPES)
        self._add([Path(p) for p in paths])

    def _add_folder(self):
        folder = filedialog.askdirectory(title="フォルダを選択")
        if not folder:
            return
        paths = [p for p in Path(folder).rglob("*") if p.suffix.lower() in self.IMG_EXTS]
        if not paths:
            messagebox.showinfo("情報", "対応する画像が見つかりませんでした。")
            return
        self._add(paths)

    def _add(self, paths):
        existing = {f["path"] for f in self.files}
        for p in paths:
            if p not in existing:
                self.files.append({"path": p, "status": "待機中", "fg": DIM})
        self._refresh()

    def _clear(self):
        if not self.converting:
            self.files.clear()
            self._refresh()

    def _choose_out(self):
        d = filedialog.askdirectory()
        if d:
            self.output_dir.set(d)

    def _show_empty(self):
        tk.Label(self.inner,
                 text="📂  画像ファイルをここに追加してください\n"
                      "HEIC / HEIF / JPEG / PNG / WebP / AVIF / BMP / TIFF",
                 font=("Yu Gothic UI", 11), bg=SURFACE,
                 fg=DIM, justify="center").pack(expand=True, pady=40)

    def _refresh(self):
        for w in self.inner.winfo_children():
            w.destroy()
        if not self.files:
            self._show_empty()
            return
        for i, f in enumerate(self.files):
            bg = SURFACE if i % 2 == 0 else ROW_ALT
            row = tk.Frame(self.inner, bg=bg)
            row.pack(fill="x")
            tk.Label(row, text=f" {f['path'].suffix.upper().lstrip('.')} ",
                     font=("Yu Gothic UI", 8, "bold"),
                     bg=BORDER, fg=MID, pady=5
                     ).pack(side="left", padx=(8, 6), pady=4)
            tk.Label(row, text=f["path"].name,
                     font=("Yu Gothic UI", 10), bg=bg, fg=TEXT, anchor="w"
                     ).pack(side="left", pady=4)
            tk.Label(row, text=f["status"],
                     font=("Yu Gothic UI", 9), bg=bg, fg=f["fg"], anchor="e"
                     ).pack(side="right", padx=10)

    def _start(self):
        if self.converting:
            return
        if not self.files:
            messagebox.showwarning("警告", "ファイルを追加してください。")
            return
        out = self.output_dir.get().strip()
        if out and not os.path.isdir(out):
            messagebox.showerror("エラー", "出力先フォルダが存在しません。")
            return
        self.converting = True
        self.prog_var.set(0)
        for f in self.files:
            f["status"] = "待機中"
            f["fg"] = DIM
        self._refresh()
        threading.Thread(target=self._run, args=(out,), daemon=True).start()

    def _run(self, out_dir):
        fmt  = self.out_format.get()
        info = self.FORMATS[fmt]
        q    = self.quality.get()
        total = len(self.files)
        ok = fail = 0
        for i, f in enumerate(self.files):
            self._set_st(i, "変換中...", WARN_C)
            try:
                dest_dir = Path(out_dir) if out_dir else f["path"].parent
                dest = dest_dir / (f["path"].stem + info["ext"])
                with Image.open(f["path"]) as img:
                    if not info["alpha"] and img.mode in ("RGBA", "LA", "P"):
                        bg = Image.new("RGB", img.size, (255, 255, 255))
                        src = img.convert("RGBA") if img.mode != "RGBA" else img
                        bg.paste(src, mask=src.split()[3])
                        img_out = bg
                    elif info["alpha"] and img.mode == "P":
                        img_out = img.convert("RGBA")
                    else:
                        img_out = img
                    if not info["alpha"] and img_out.mode not in ("RGB", "L"):
                        img_out = img_out.convert("RGB")
                    elif info["alpha"] and img_out.mode not in ("RGB", "RGBA", "L", "LA"):
                        img_out = img_out.convert("RGBA")
                    kw = {}
                    if info["quality"]:
                        kw["quality"] = q
                    if fmt == "JPEG":
                        kw["optimize"] = True
                    img_out.save(dest, info["pil"], **kw)
                self._set_st(i, "✅ 完了", SUCCESS)
                ok += 1
            except Exception:
                self._set_st(i, "❌ 失敗", ERROR_C)
                fail += 1
            self.after(0, lambda p=(i+1)/total*100, c=i+1, t=total:
                       (self.prog_var.set(p),
                        self.prog_label.config(text=f"変換中... {c} / {t}")))
        self.after(0, lambda: self._done(ok, fail))

    def _set_st(self, i, s, fg):
        self.files[i]["status"] = s
        self.files[i]["fg"] = fg
        self.after(0, self._refresh)

    def _done(self, ok, fail):
        self.converting = False
        self.prog_var.set(100)
        self.prog_label.config(
            text=f"🎉 完了：{ok} 件成功" + (f"  /  {fail} 件失敗" if fail else ""),
            fg=SUCCESS if not fail else WARN_C)
        if fail == 0:
            messagebox.showinfo("完了", f"🎉 {ok} 件の変換が完了しました！")
        else:
            messagebox.showwarning("完了", f"成功: {ok} 件 / 失敗: {fail} 件")


# ════════════════════════════════════════════════════════════════════════════
# 動画変換タブ
# ════════════════════════════════════════════════════════════════════════════
class VideoTab(tk.Frame):
    FORMATS = {
        "MP4":  {"ext": ".mp4",  "note": "最も互換性が高い・スマホ・PC問わず再生可能"},
        "WebM": {"ext": ".webm", "note": "Web向け・軽量・ブラウザ再生に最適"},
        "GIF":  {"ext": ".gif",  "note": "アニメーションGIF（音声なし）"},
        "MP3":  {"ext": ".mp3",  "note": "音声のみ抽出して保存"},
    }
    FILETYPES = [("対応動画", "*.mp4 *.mov *.avi *.mkv *.webm *.m4v"), ("すべて", "*.*")]

    def __init__(self, parent):
        super().__init__(parent, bg=BG)
        self.filepath   = None
        self.converting = False
        self.out_format = tk.StringVar(value="MP4")
        self.gif_fps    = tk.IntVar(value=10)
        self.gif_width  = tk.IntVar(value=480)
        self._build()

    def _build(self):
        if not ffmpeg_ok():
            o, i = make_card(self)
            o.pack(padx=20, pady=30, fill="x")
            tk.Label(i, text="⚠️  FFmpegがインストールされていないため\n動画変換は使用できません。",
                     font=("Yu Gothic UI", 12), bg=SURFACE, fg=ERROR_C,
                     justify="center", pady=16).pack()
            return

        build_guide(self,
            ["① 動画ファイルを選択 でファイルを選ぶ",
             "② 変換後の形式を選ぶ（GIFはサイズも設定可）",
             "③「✨ 変換開始」ボタンを押す"],
            "MP4 / MOV / AVI / MKV / WebM / M4V")

        so, si = make_card(self)
        so.pack(fill="x", padx=16, pady=(10, 0))
        sr = tk.Frame(si, bg=SURFACE)
        sr.pack(fill="x", padx=12, pady=10)
        SecBtn(sr, "📂 動画ファイルを選択", self._pick).pack(side="left")
        self.file_label = tk.Label(sr, text="未選択",
                                   font=("Yu Gothic UI", 10),
                                   bg=SURFACE, fg=DIM, anchor="w")
        self.file_label.pack(side="left", padx=(12, 0))

        co, ci = make_card(self)
        co.pack(fill="x", padx=16, pady=(8, 0))

        r1 = tk.Frame(ci, bg=SURFACE)
        r1.pack(fill="x", padx=12, pady=(10, 4))
        tk.Label(r1, text="✨ 変換後の形式", font=("Yu Gothic UI", 10, "bold"),
                 bg=SURFACE, fg=MID).pack(side="left")
        cb = build_cb(r1, self.out_format, list(self.FORMATS.keys()))
        cb.pack(side="left", padx=(8, 0))
        cb.bind("<<ComboboxSelected>>", self._on_fmt)
        self.fmt_note = tk.Label(r1, text="", font=("Yu Gothic UI", 9),
                                 bg=SURFACE, fg=MID)
        self.fmt_note.pack(side="left", padx=(10, 0))

        self.gif_frame = tk.Frame(ci, bg=SURFACE)
        for lbl, var, opts in [("🎞️ fps", self.gif_fps, [5, 10, 15, 24]),
                                ("📐 横幅", self.gif_width, [240, 360, 480, 640])]:
            row = tk.Frame(self.gif_frame, bg=SURFACE)
            row.pack(fill="x", padx=12, pady=(0, 2))
            tk.Label(row, text=lbl, font=("Yu Gothic UI", 9),
                     bg=SURFACE, fg=DIM, width=7, anchor="w").pack(side="left")
            for v in opts:
                tk.Radiobutton(row, text=str(v) if lbl.endswith("fps") else f"{v}px",
                               variable=var, value=v,
                               bg=SURFACE, fg=TEXT, selectcolor=CARD,
                               activebackground=SURFACE,
                               font=("Yu Gothic UI", 9)).pack(side="left", padx=4)

        r2 = tk.Frame(ci, bg=SURFACE)
        r2.pack(fill="x", padx=12, pady=(0, 10))
        tk.Label(r2, text="出力先", font=("Yu Gothic UI", 9),
                 bg=SURFACE, fg=DIM, width=7, anchor="w").pack(side="left")
        self.out_dir = tk.StringVar(value="")
        tk.Entry(r2, textvariable=self.out_dir,
                 bg=CARD, fg=TEXT, insertbackground=MID,
                 relief="flat", font=("Yu Gothic UI", 9), bd=3,
                 highlightbackground=BORDER, highlightthickness=1).pack(
            side="left", fill="x", expand=True, padx=(4, 8))
        SecBtn(r2, "選択", self._choose_out).pack(side="left")

        lo, li = make_card(self)
        lo.pack(fill="both", expand=True, padx=16, pady=(8, 0))
        tk.Label(li, text="📝 変換ログ",
                 font=("Yu Gothic UI", 9, "bold"),
                 bg=SURFACE, fg=MID,
                 anchor="w", padx=10, pady=4).pack(fill="x")
        ll = tk.Frame(li, bg=SURFACE)
        ll.pack(fill="both", expand=True)
        self.log = tk.Text(ll, bg=CARD, fg=TEXT, font=("Consolas", 9),
                           relief="flat", state="disabled", wrap="word",
                           highlightthickness=0)
        sb = ttk.Scrollbar(ll, command=self.log.yview)
        self.log.configure(yscrollcommand=sb.set)
        sb.pack(side="right", fill="y")
        self.log.pack(fill="both", expand=True, padx=6, pady=(0, 6))

        PillBtn(self, "✨ 変換開始", self._start).pack(
            fill="x", padx=16, pady=(8, 16))
        self._on_fmt()

    def _on_fmt(self, _=None):
        fmt = self.out_format.get()
        self.fmt_note.config(text=f"📌 {self.FORMATS[fmt]['note']}")
        if fmt == "GIF":
            self.gif_frame.pack(fill="x")
        else:
            self.gif_frame.pack_forget()

    def _pick(self):
        path = filedialog.askopenfilename(title="動画ファイルを選択", filetypes=self.FILETYPES)
        if path:
            self.filepath = Path(path)
            mb = round(self.filepath.stat().st_size / 1024 / 1024, 1)
            self.file_label.config(text=f"📄 {self.filepath.name}  ({mb} MB)", fg=TEXT)

    def _choose_out(self):
        d = filedialog.askdirectory()
        if d:
            self.out_dir.set(d)

    def _log(self, msg):
        self.log.configure(state="normal")
        self.log.insert("end", msg + "\n")
        self.log.see("end")
        self.log.configure(state="disabled")

    def _start(self):
        if self.converting:
            return
        if not self.filepath or not self.filepath.exists():
            messagebox.showwarning("警告", "動画ファイルを選択してください。")
            return
        self.converting = True
        self.log.configure(state="normal")
        self.log.delete("1.0", "end")
        self.log.configure(state="disabled")
        threading.Thread(target=self._run, daemon=True).start()

    def _run(self):
        import shutil
        fmt  = self.out_format.get()
        ext  = self.FORMATS[fmt]["ext"]
        out_d = self.out_dir.get().strip()
        dest_dir = Path(out_d) if out_d else self.filepath.parent
        out_name = self.filepath.stem + ext
        out_path = dest_dir / out_name
        try:
            with tempfile.TemporaryDirectory() as tmp:
                inp = os.path.join(tmp, "input" + self.filepath.suffix)
                shutil.copy2(self.filepath, inp)
                if fmt == "MP4":
                    cmd = ["ffmpeg", "-i", inp, "-c:v", "libx264",
                           "-c:a", "aac", "-movflags", "+faststart", "-y", str(out_path)]
                elif fmt == "WebM":
                    cmd = ["ffmpeg", "-i", inp, "-c:v", "libvpx-vp9",
                           "-c:a", "libopus", "-y", str(out_path)]
                elif fmt == "GIF":
                    fps, width = self.gif_fps.get(), self.gif_width.get()
                    pal = os.path.join(tmp, "palette.png")
                    subprocess.run(["ffmpeg", "-i", inp,
                                    "-vf", f"fps={fps},scale={width}:-1:flags=lanczos,palettegen",
                                    "-y", pal],
                                   capture_output=True, timeout=120, creationflags=_NO_WIN)
                    cmd = ["ffmpeg", "-i", inp, "-i", pal,
                           "-filter_complex",
                           f"fps={fps},scale={width}:-1:flags=lanczos[x];[x][1:v]paletteuse",
                           "-y", str(out_path)]
                elif fmt == "MP3":
                    cmd = ["ffmpeg", "-i", inp, "-q:a", "2", "-map", "a", "-y", str(out_path)]

                self.after(0, lambda: self._log(f"変換開始: {self.filepath.name} → {out_name}"))
                result = subprocess.run(cmd, capture_output=True, timeout=600, creationflags=_NO_WIN)
                if result.returncode == 0:
                    mb = round(out_path.stat().st_size / 1024 / 1024, 1)
                    self.after(0, lambda: self._log(f"✅ 完了: {out_name}  ({mb} MB)"))
                    self.after(0, lambda: messagebox.showinfo(
                        "完了", f"🎉 変換が完了しました！\n保存先: {out_path}"))
                else:
                    err = result.stderr.decode("utf-8", errors="replace")[-300:]
                    self.after(0, lambda: self._log(f"❌ 失敗:\n{err}"))
                    self.after(0, lambda: messagebox.showerror("失敗", "変換に失敗しました。"))
        except subprocess.TimeoutExpired:
            self.after(0, lambda: messagebox.showerror("タイムアウト", "変換がタイムアウトしました。"))
        except Exception as e:
            self.after(0, lambda e=e: messagebox.showerror("エラー", str(e)))
        finally:
            self.after(0, lambda: setattr(self, "converting", False))


# ════════════════════════════════════════════════════════════════════════════
# 文書変換タブ
# ════════════════════════════════════════════════════════════════════════════
class DocTab(tk.Frame):
    OUTPUT_BY_EXT = {
        "txt":  ["PDF", "DOCX", "JPEG", "PNG", "WebP"],
        "docx": ["PDF", "TXT",  "JPEG", "PNG", "WebP"],
        "pdf":  ["TXT", "DOCX", "PPTX", "JPEG", "PNG", "WebP"],
        "xlsx": ["PDF", "CSV",  "JPEG", "PNG", "WebP"],
        "pptx": ["PDF", "TXT",  "DOCX", "JPEG", "PNG", "WebP"],
    }
    FORMAT_NOTES = {
        "PDF":  "どの端末でも同じレイアウトで開ける",
        "TXT":  "テキストを抽出してプレーンテキストで保存",
        "CSV":  "Excel・スプレッドシートで開けるデータ形式",
        "DOCX": "Word文書として開ける（スキャン画像PDFはOCRで文字認識）",
        "PPTX": "PDFの各ページをスライドとしてテキスト化（スキャン画像PDFはOCR）",
        "JPEG": "ページごとに画像化（JPEG）",
        "PNG":  "ページごとに画像化（PNG）",
        "WebP": "ページごとに画像化（WebP）",
    }
    FILETYPES = [("対応文書", "*.txt *.docx *.pdf *.xlsx *.pptx"), ("すべて", "*.*")]

    def __init__(self, parent):
        super().__init__(parent, bg=BG)
        self.filepath   = None
        self.converting = False
        self.out_format = tk.StringVar(value="")
        self.dpi        = tk.IntVar(value=150)
        self._build()

    def _build(self):
        build_guide(self,
            ["① 文書ファイルを選択 でファイルを選ぶ",
             "② 変換後の形式を選ぶ（ファイルの種類で選択肢が変わります）",
             "③「✨ 変換開始」ボタンを押す"],
            "TXT / DOCX（Word） / PDF / XLSX（Excel） / PPTX（PowerPoint）")

        so, si = make_card(self)
        so.pack(fill="x", padx=16, pady=(10, 0))
        sr = tk.Frame(si, bg=SURFACE)
        sr.pack(fill="x", padx=12, pady=10)
        SecBtn(sr, "📂 文書ファイルを選択", self._pick).pack(side="left")
        self.file_label = tk.Label(sr, text="未選択",
                                   font=("Yu Gothic UI", 10),
                                   bg=SURFACE, fg=DIM, anchor="w")
        self.file_label.pack(side="left", padx=(12, 0))

        co, ci = make_card(self)
        co.pack(fill="x", padx=16, pady=(8, 0))

        r1 = tk.Frame(ci, bg=SURFACE)
        r1.pack(fill="x", padx=12, pady=(10, 4))
        tk.Label(r1, text="✨ 変換後の形式", font=("Yu Gothic UI", 10, "bold"),
                 bg=SURFACE, fg=MID).pack(side="left")
        self.fmt_cb = build_cb(r1, self.out_format, [])
        self.fmt_cb.pack(side="left", padx=(8, 0))
        self.fmt_cb.bind("<<ComboboxSelected>>", self._on_fmt)
        self.fmt_note = tk.Label(r1, text="ファイルを選択してください",
                                 font=("Yu Gothic UI", 9), bg=SURFACE, fg=DIM)
        self.fmt_note.pack(side="left", padx=(10, 0))

        self.dpi_frame = tk.Frame(ci, bg=SURFACE)
        df = tk.Frame(self.dpi_frame, bg=SURFACE)
        df.pack(fill="x", padx=12, pady=(0, 4))
        tk.Label(df, text="🖨️ 解像度", font=("Yu Gothic UI", 9),
                 bg=SURFACE, fg=DIM, width=7, anchor="w").pack(side="left")
        for v, lbl in [(96, "96 dpi（画面）"), (150, "150 dpi（標準）"), (300, "300 dpi（印刷）")]:
            tk.Radiobutton(df, text=lbl, variable=self.dpi, value=v,
                           bg=SURFACE, fg=TEXT, selectcolor=CARD,
                           activebackground=SURFACE,
                           font=("Yu Gothic UI", 9)).pack(side="left", padx=6)

        r2 = tk.Frame(ci, bg=SURFACE)
        r2.pack(fill="x", padx=12, pady=(0, 10))
        tk.Label(r2, text="出力先", font=("Yu Gothic UI", 9),
                 bg=SURFACE, fg=DIM, width=7, anchor="w").pack(side="left")
        self.out_dir = tk.StringVar(value="")
        tk.Entry(r2, textvariable=self.out_dir,
                 bg=CARD, fg=TEXT, insertbackground=MID,
                 relief="flat", font=("Yu Gothic UI", 9), bd=3,
                 highlightbackground=BORDER, highlightthickness=1).pack(
            side="left", fill="x", expand=True, padx=(4, 8))
        SecBtn(r2, "選択", self._choose_out).pack(side="left")

        lo, li = make_card(self)
        lo.pack(fill="both", expand=True, padx=16, pady=(8, 0))
        tk.Label(li, text="📝 変換ログ",
                 font=("Yu Gothic UI", 9, "bold"),
                 bg=SURFACE, fg=MID, anchor="w", padx=10, pady=4).pack(fill="x")
        ll = tk.Frame(li, bg=SURFACE)
        ll.pack(fill="both", expand=True)
        self.log = tk.Text(ll, bg=CARD, fg=TEXT, font=("Consolas", 9),
                           relief="flat", state="disabled", wrap="word",
                           highlightthickness=0)
        sb = ttk.Scrollbar(ll, command=self.log.yview)
        self.log.configure(yscrollcommand=sb.set)
        sb.pack(side="right", fill="y")
        self.log.pack(fill="both", expand=True, padx=6, pady=(0, 6))

        PillBtn(self, "✨ 変換開始", self._start).pack(
            fill="x", padx=16, pady=(8, 16))

    def _pick(self):
        path = filedialog.askopenfilename(title="文書ファイルを選択", filetypes=self.FILETYPES)
        if not path:
            return
        self.filepath = Path(path)
        kb = round(self.filepath.stat().st_size / 1024, 1)
        self.file_label.config(text=f"📄 {self.filepath.name}  ({kb} KB)", fg=TEXT)
        ext  = self.filepath.suffix.lower().lstrip(".")
        opts = self.OUTPUT_BY_EXT.get(ext, [])
        self.fmt_cb.configure(values=opts)
        if opts:
            self.out_format.set(opts[0])
            self._on_fmt()
        else:
            self.out_format.set("")
            messagebox.showwarning("非対応", "この形式は対応していません。")

    def _on_fmt(self, _=None):
        fmt = self.out_format.get()
        note = self.FORMAT_NOTES.get(fmt, "")
        self.fmt_note.config(text=f"📌 {note}" if note else "", fg=MID)
        if fmt in ("JPEG", "PNG", "WebP"):
            self.dpi_frame.pack(fill="x")
        else:
            self.dpi_frame.pack_forget()

    def _choose_out(self):
        d = filedialog.askdirectory()
        if d:
            self.out_dir.set(d)

    def _log(self, msg):
        self.log.configure(state="normal")
        self.log.insert("end", msg + "\n")
        self.log.see("end")
        self.log.configure(state="disabled")

    def _start(self):
        if self.converting:
            return
        if not self.filepath or not self.filepath.exists():
            messagebox.showwarning("警告", "ファイルを選択してください。")
            return
        if not self.out_format.get():
            messagebox.showwarning("警告", "変換形式を選択してください。")
            return
        self.converting = True
        self.log.configure(state="normal")
        self.log.delete("1.0", "end")
        self.log.configure(state="disabled")
        threading.Thread(target=self._run, args=(self.out_format.get(),), daemon=True).start()

    def _run(self, fmt):
        import shutil
        try:
            ext      = self.filepath.suffix.lower().lstrip(".")
            out_d    = self.out_dir.get().strip()
            dest_dir = Path(out_d) if out_d else self.filepath.parent
            stem     = self.filepath.stem
            dpi      = self.dpi.get()
            self.after(0, lambda: self._log(f"変換開始: {self.filepath.name} → {fmt}"))
            with tempfile.TemporaryDirectory() as tmp:
                inp = os.path.join(tmp, "input." + ext)
                shutil.copy2(self.filepath, inp)
                if fmt == "PDF":
                    lo = libreoffice_ok()
                    if not lo:
                        self.after(0, lambda: messagebox.showerror(
                            "エラー", "LibreOfficeがインストールされていないためPDF変換できません。"))
                        return
                    env = {**os.environ, "HOME": tmp}
                    r = subprocess.run([lo, "--headless", "--norestore",
                                        "--convert-to", "pdf",
                                        "--outdir", str(dest_dir), inp],
                                       capture_output=True, timeout=120,
                                       env=env, creationflags=_NO_WIN)
                    op = dest_dir / (stem + ".pdf")
                    ok = r.returncode == 0
                    self.after(0, lambda ok=ok, op=op: (
                        self._log(f"{'✅ 完了' if ok else '❌ 失敗'}: {op.name}"),
                        messagebox.showinfo("完了", f"🎉 保存先: {op}") if ok
                        else messagebox.showerror("失敗", "PDF変換に失敗しました。")))
                elif fmt == "TXT":
                    text = self._extract_text(inp, ext)
                    op = dest_dir / (stem + ".txt")
                    op.write_text(text, encoding="utf-8")
                    self.after(0, lambda op=op: (self._log(f"✅ 完了: {op.name}"),
                                                  messagebox.showinfo("完了", f"🎉 保存先: {op}")))
                elif fmt == "CSV":
                    import openpyxl
                    wb = openpyxl.load_workbook(inp, data_only=True)
                    ws = wb.active
                    op = dest_dir / (stem + ".csv")
                    with open(op, "w", newline="", encoding="utf-8-sig") as f:
                        w = csv.writer(f)
                        for row in ws.iter_rows(values_only=True):
                            w.writerow(['' if v is None else v for v in row])
                    self.after(0, lambda op=op: (self._log(f"✅ 完了: {op.name}"),
                                                  messagebox.showinfo("完了", f"🎉 保存先: {op}")))
                elif fmt == "DOCX":
                    op = dest_dir / (stem + ".docx")
                    if ext == "pdf":
                        if self._pdf_is_scanned(inp):
                            if not tesseract_ok():
                                self.after(0, lambda: messagebox.showerror(
                                    "エラー",
                                    "このPDFは文字情報を持たないスキャン画像PDFのため、"
                                    "OCR（文字認識）が必要です。\n"
                                    "Tesseract-OCRがインストールされていないため変換できません。"))
                                return
                            self.after(0, lambda: self._log(
                                "📷 スキャン画像PDFを検出しました。OCRで文字認識します..."))
                            self._pdf_to_docx_ocr(inp, op)
                        else:
                            from pdf2docx import Converter
                            cv = Converter(inp); cv.convert(str(op)); cv.close()
                    elif ext == "pptx":
                        from pptx import Presentation
                        from docx import Document
                        prs = Presentation(inp)
                        doc = Document()
                        for i, slide in enumerate(prs.slides, 1):
                            doc.add_heading(f"スライド {i}", level=1)
                            for shape in slide.shapes:
                                if shape.has_text_frame:
                                    for para in shape.text_frame.paragraphs:
                                        if para.text.strip():
                                            doc.add_paragraph(para.text.strip())
                            doc.add_paragraph("")
                        doc.save(str(op))
                    else:
                        from docx import Document
                        doc = Document()
                        with open(inp, "r", encoding="utf-8", errors="replace") as f:
                            for line in f.read().splitlines():
                                doc.add_paragraph(line)
                        doc.save(str(op))
                    self.after(0, lambda op=op: (self._log(f"✅ 完了: {op.name}"),
                                                  messagebox.showinfo("完了", f"🎉 保存先: {op}")))
                elif fmt == "PPTX":
                    op = dest_dir / (stem + ".pptx")
                    if self._pdf_is_scanned(inp) and not tesseract_ok():
                        self.after(0, lambda: messagebox.showerror(
                            "エラー",
                            "このPDFは文字情報を持たないスキャン画像PDFのため、"
                            "OCR（文字認識）が必要です。\n"
                            "Tesseract-OCRがインストールされていないため変換できません。"))
                        return
                    self._pdf_to_pptx(inp, op)
                    self.after(0, lambda op=op: (self._log(f"✅ 完了: {op.name}"),
                                                  messagebox.showinfo("完了", f"🎉 保存先: {op}")))
                elif fmt in ("JPEG", "PNG", "WebP"):
                    images = self._to_images(inp, ext, fmt, dpi, tmp)
                    ext_map = {"JPEG": ".jpg", "PNG": ".png", "WebP": ".webp"}
                    if len(images) == 1:
                        op = dest_dir / (stem + ext_map[fmt])
                        op.write_bytes(images[0][1])
                        self.after(0, lambda op=op: (self._log(f"✅ 完了: {op.name}"),
                                                      messagebox.showinfo("完了", f"🎉 保存先: {op}")))
                    else:
                        zp = dest_dir / (stem + "_images.zip")
                        with zipfile.ZipFile(zp, "w", zipfile.ZIP_DEFLATED) as zf:
                            for name, data in images:
                                zf.writestr(name, data)
                        self.after(0, lambda zp=zp, n=len(images): (
                            self._log(f"✅ 完了: {n} ページ"),
                            messagebox.showinfo("完了", f"🎉 {n} ページ保存しました！\n保存先: {zp}")))
        except subprocess.TimeoutExpired:
            self.after(0, lambda: messagebox.showerror("タイムアウト", "変換がタイムアウトしました。"))
        except Exception as e:
            self.after(0, lambda e=e: (self._log(f"❌ エラー: {e}"),
                                       messagebox.showerror("エラー", str(e))))
        finally:
            self.after(0, lambda: setattr(self, "converting", False))

    def _extract_text(self, path, ext):
        if ext == "pdf":
            if self._pdf_is_scanned(path):
                if not tesseract_ok():
                    raise RuntimeError(
                        "このPDFは文字情報を持たないスキャン画像PDFのため、OCR（文字認識）が必要です。"
                        "Tesseract-OCRがインストールされていないため変換できません。")
                self.after(0, lambda: self._log(
                    "📷 スキャン画像PDFを検出しました。OCRで文字認識します..."))
                return self._pdf_ocr_text(path)
            import pdfplumber
            parts = []
            with pdfplumber.open(path) as pdf:
                for page in pdf.pages:
                    t = page.extract_text()
                    if t:
                        parts.append(t)
            return "\n\n".join(parts)
        elif ext == "docx":
            from docx import Document
            return "\n".join(p.text for p in Document(path).paragraphs)
        elif ext == "pptx":
            from pptx import Presentation
            lines = []
            for i, slide in enumerate(Presentation(path).slides, 1):
                lines.append(f"=== スライド {i} ===")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        lines.append(shape.text_frame.text)
                lines.append("")
            return "\n".join(lines)
        return ""

    def _pdf_is_scanned(self, inp, threshold=20):
        import fitz
        doc = fitz.open(inp)
        n = min(len(doc), 3)
        total_chars = sum(len(doc[i].get_text().strip()) for i in range(n))
        doc.close()
        return (total_chars / max(n, 1)) < threshold

    def _ocr_page_text(self, exe, img, langs_available):
        bundled = bundled_tessdata_dir()
        vert_dir = bundled if bundled and os.path.isfile(
            os.path.join(bundled, "jpn_vert.traineddata")) else None
        candidates = []
        conf_h = _run_tesseract(exe, img, "jpn", want_conf=True)
        candidates.append((conf_h, "jpn", None, None))
        if "jpn_vert" in langs_available:
            conf_v = _run_tesseract(exe, img, "jpn_vert", tessdata_dir=vert_dir, psm=5, want_conf=True)
            candidates.append((conf_v, "jpn_vert", vert_dir, 5))
        candidates.sort(key=lambda c: c[0], reverse=True)
        _, lang, tessdata_dir, psm = candidates[0]
        return _run_tesseract(exe, img, lang, tessdata_dir=tessdata_dir, psm=psm)

    def _iter_ocr_pages(self, inp):
        import fitz
        exe = tesseract_path()
        langs_available = tesseract_langs()
        if "jpn_vert" not in langs_available:
            self.after(0, lambda: self._log(
                "⚠️ jpn_vert（縦書き用言語データ）が未インストールのため、"
                "縦書き文書は認識精度が下がる場合があります。"))
        doc = fitz.open(inp)
        matrix = fitz.Matrix(300 / 72, 300 / 72)
        total = len(doc)
        for i, page in enumerate(doc):
            self.after(0, lambda i=i, t=total: self._log(f"OCR処理中... {i+1} / {t} ページ"))
            pix = page.get_pixmap(matrix=matrix, alpha=False)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            text = self._ocr_page_text(exe, img, langs_available)
            lines = [clean_ocr_ja(ln.strip()) for ln in text.splitlines() if ln.strip()]
            yield i, total, lines
        doc.close()

    def _pdf_ocr_text(self, inp):
        return "\n\n".join("\n".join(lines) for _, _, lines in self._iter_ocr_pages(inp))

    def _pdf_to_docx_ocr(self, inp, out_path):
        from docx import Document
        wdoc = Document()
        for i, total, lines in self._iter_ocr_pages(inp):
            if lines:
                for line in lines:
                    wdoc.add_paragraph(line)
            else:
                wdoc.add_paragraph("")
            if i < total - 1:
                wdoc.add_page_break()
        wdoc.save(str(out_path))

    def _pdf_to_pptx(self, inp, out_path):
        import fitz
        from pptx import Presentation
        from pptx.util import Inches, Pt
        doc = fitz.open(inp)
        prs = Presentation()
        if len(doc) > 0:
            r = doc[0].rect
            prs.slide_width  = Inches(r.width  / 72)
            prs.slide_height = Inches(r.height / 72)
        total = len(doc)
        doc.close()

        if self._pdf_is_scanned(inp):
            self.after(0, lambda: self._log(
                "📷 スキャン画像PDFを検出しました。OCRで文字認識します..."))
            page_texts = ["\n".join(lines) for _, _, lines in self._iter_ocr_pages(inp)]
        else:
            import pdfplumber
            with pdfplumber.open(inp) as pdf:
                page_texts = [(p.extract_text() or "") for p in pdf.pages]

        blank = prs.slide_layouts[6]
        margin = Inches(0.4)
        for i in range(total):
            slide = prs.slides.add_slide(blank)
            tb = slide.shapes.add_textbox(
                margin, margin, prs.slide_width - margin * 2, prs.slide_height - margin * 2)
            tf = tb.text_frame
            tf.word_wrap = True
            lines = [ln for ln in page_texts[i].splitlines() if ln.strip()] or [""]
            tf.text = lines[0]
            tf.paragraphs[0].font.size = Pt(18)
            for ln in lines[1:]:
                p = tf.add_paragraph()
                p.text = ln
                p.font.size = Pt(18)
        prs.save(str(out_path))

    def _to_images(self, inp, ext, fmt, dpi, tmp):
        lo = libreoffice_ok()
        if ext != "pdf":
            if not lo:
                raise RuntimeError("LibreOfficeがインストールされていません。")
            env = {**os.environ, "HOME": tmp}
            subprocess.run([lo, "--headless", "--norestore",
                            "--convert-to", "pdf", "--outdir", tmp, inp],
                           capture_output=True, timeout=120,
                           env=env, creationflags=_NO_WIN)
            inp = os.path.join(tmp, Path(inp).stem + ".pdf")
        import fitz
        doc    = fitz.open(inp)
        matrix = fitz.Matrix(dpi / 72, dpi / 72)
        suffix  = {"JPEG": ".jpg", "PNG": ".png", "WebP": ".webp"}[fmt]
        pil_fmt = {"JPEG": "JPEG", "PNG": "PNG", "WebP": "WEBP"}[fmt]
        images  = []
        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=matrix, alpha=False)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            buf = io.BytesIO()
            kw  = {"quality": 90} if fmt in ("JPEG", "WebP") else {}
            img.save(buf, pil_fmt, **kw)
            buf.seek(0)
            images.append((f"page_{i+1:03d}{suffix}", buf.getvalue()))
        doc.close()
        return images


# ════════════════════════════════════════════════════════════════════════════
# メインアプリ
# ════════════════════════════════════════════════════════════════════════════
class App(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("拡張子変換ツール")
        self.geometry("780x690")
        self.minsize(660, 580)
        self.configure(bg=BG)
        self.resizable(True, True)
        self._build()

    def _build(self):
        # ヘッダー
        hdr = tk.Frame(self, bg=BG)
        hdr.pack(fill="x", pady=(16, 8))
        tk.Label(hdr, text="🔄", font=("Yu Gothic UI", 22),
                 bg=BG, fg=MID).pack()
        tk.Label(hdr, text="拡張子変換ツール",
                 font=("Yu Gothic UI", 20, "bold"),
                 bg=BG, fg=MID).pack()
        tk.Label(hdr, text="ドラッグ＆ドロップで簡単変換 ✨",
                 font=("Yu Gothic UI", 9),
                 bg=BG, fg=DIM).pack()

        # タブバー（ウェブ版と同じ丸みのある紫バー）
        tab_wrap = tk.Frame(self, bg=BG)
        tab_wrap.pack(fill="x", padx=16, pady=(4, 0))
        tab_bar = tk.Frame(tab_wrap, bg=TAB_BAR,
                           highlightbackground=BORDER, highlightthickness=1)
        tab_bar.pack(fill="x", padx=0, pady=0)

        self._tab_btns = []
        for i, (label, key) in enumerate([
            ("　🖼️  画像変換　", "img"),
            ("　🎬  動画変換　", "vid"),
            ("　📄  文書変換　", "doc"),
        ]):
            btn = tk.Button(
                tab_bar, text=label,
                font=("Yu Gothic UI", 11, "bold"),
                relief="flat", bd=0, cursor="hand2",
                pady=10,
                command=lambda i=i: self._show(i))
            btn.pack(side="left", fill="both", expand=True)
            self._tab_btns.append(btn)

        # コンテンツ
        self._tabs = [ImageTab(self), VideoTab(self), DocTab(self)]
        for tab in self._tabs:
            tab.pack(fill="both", expand=True)
            tab.pack_forget()

        self._show(0)

    def _show(self, index):
        for i, btn in enumerate(self._tab_btns):
            if i == index:
                btn.config(bg=TAB_ACT, fg="#ffffff",
                           activebackground=BTN_HOV, activeforeground="#ffffff")
            else:
                btn.config(bg=TAB_BAR, fg=MID,
                           activebackground=TAB_HOV, activeforeground=MID)
        for i, tab in enumerate(self._tabs):
            tab.pack(fill="both", expand=True) if i == index else tab.pack_forget()


if __name__ == "__main__":
    App().mainloop()
